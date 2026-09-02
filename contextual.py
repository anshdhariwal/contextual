#!/usr/bin/env python3
"""
contextual — developed by @anshdhariwal
Run with: python contextual.py
"""

import sys, os, re, time, subprocess, io, hashlib

for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8", errors="replace")

def _importable(pkg):
    try: __import__(pkg); return True
    except ImportError: return False

_needed = {"pdfplumber": "pdfplumber", "pypdf": "pypdf",
           "docx": "python-docx", "pptx": "python-pptx", "rich": "rich"}
_missing = [pip_name for mod, pip_name in _needed.items() if not _importable(mod)]
if _missing:
    print(f"\n  Installing: {', '.join(_missing)} …\n")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "--quiet"] + _missing)

import pdfplumber
from pypdf import PdfReader
import docx as _docx
from docx.oxml.ns import qn as _qn
from pptx import Presentation as _Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from rich.console import Console
from rich.panel import Panel
from rich.table import Table
from rich.text import Text
from rich.align import Align
from rich import box
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TaskProgressColumn

console = Console()
FAST_MODE = "-f" in sys.argv or "--fast" in sys.argv
_STEP_OFFSET = 0
_TOTAL_STEPS = 5

OCR_MIN_BYTES = 1024
OCR_MIN_SIDE = 64
OCR_MAX_SIDE = 2000

_ocr = False
_ocr_cache = {}

def get_ocr():
    global _ocr
    if _ocr is not False:
        return _ocr
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        print("\n  Installing OCR engine (rapidocr-onnxruntime, one-time) …")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install",
                                   "--quiet", "rapidocr-onnxruntime"])
            from rapidocr_onnxruntime import RapidOCR
        except Exception:
            print("  couldn't set up OCR — continuing without it.\n")
            _ocr = None
            return None
    print("  loading OCR model …")
    _ocr = RapidOCR()
    return _ocr

def ocr_blob(blob: bytes) -> str:
    key = hashlib.md5(blob).hexdigest()
    if key in _ocr_cache:
        return _ocr_cache[key]
    text = ""
    img = None
    try:
        img = Image.open(io.BytesIO(blob)).convert("RGB")
        w, h = img.size
        if len(blob) < OCR_MIN_BYTES or w < OCR_MIN_SIDE or h < OCR_MIN_SIDE:
            img = None
    except Exception:
        img = None
    if img is not None:
        engine = get_ocr()
        if engine is not None:
            try:
                if max(w, h) > OCR_MAX_SIDE:
                    scale = OCR_MAX_SIDE / max(w, h)
                    img = img.resize((int(w * scale), int(h * scale)))
                import numpy as np
                result, _ = engine(np.array(img)[:, :, ::-1])
                if result:
                    result.sort(key=lambda r: r[0][0][1])
                    text = "\n".join(t for _, t, s in result)
            except Exception:
                text = ""
    _ocr_cache[key] = text
    return text

def W() -> int:
    return min(console.width, 90)

def clear():
    os.system("cls" if os.name == "nt" else "clear")

def measure_height(renderable) -> int:
    buf = io.StringIO()
    tmp = Console(file=buf, width=W(), highlight=False, no_color=True)
    tmp.print(renderable)
    return buf.getvalue().count("\n")

def redraw_table(t, prev_h: int) -> int:
    if prev_h > 0:
        sys.stdout.write(f"\033[{prev_h}A\033[J")
        sys.stdout.flush()
    console.print(t)
    return measure_height(t)

def banner():
    w = W()
    fast = FAST_MODE
    acc = "red" if fast else "green"
    title = Text(justify="center")
    title.append("◈  ", style=f"dim {acc}")
    title.append("C O N T E X T U A L", style=f"bold {acc}")
    title.append("  ◈", style=f"dim {acc}")
    deco = Text("─ ─ ─ ─ ─ ─ ─ ─ ─ ─", justify="center", style=f"dim {acc}")
    sub = Text(justify="center")
    sub.append("developed by ", style="dim white")
    sub.append("@anshdhariwal", style="bold magenta")
    if fast:
        sub.append("  ·  ", style="dim red")
        sub.append("FAST MODE", style="bold red")
    content = Align.center(Text.assemble(title, "\n", deco, "\n", sub))
    console.print(
        Panel(content, box=box.DOUBLE_EDGE, border_style=acc,
              width=w, padding=(1, 10))
    )
    console.print()

def numeric_key(name: str):
    m = re.match(r"^(\d+)", name)
    return (int(m.group(1)), name) if m else (float("inf"), name)

def discover_files() -> tuple:
    all_files = os.listdir(".")
    pdfs = sorted([f for f in all_files if f.lower().endswith(".pdf")], key=numeric_key)
    docxs = sorted([f for f in all_files if f.lower().endswith(".docx")], key=numeric_key)
    pptxs = sorted([f for f in all_files if f.lower().endswith(".pptx")], key=numeric_key)
    return pdfs, docxs, pptxs

def make_table(files, highlight=None, ghost=None, arrived=None):
    w = W()
    t = Table(
        show_header=True, header_style="bold magenta",
        box=box.SIMPLE_HEAVY, border_style="dim",
        width=w - 4, pad_edge=True,
    )
    t.add_column("#", justify="right", width=4, no_wrap=True)
    t.add_column("Filename", style="white")
    cap = w - 20
    for i, f in enumerate(files, 1):
        label = (f[: cap - 1] + "…") if len(f) > cap else f
        if arrived and i == arrived:
            t.add_row(f"[bold green]{i}[/bold green]",
                      f"[bold green]{label}[/bold green]")
        elif ghost and i == ghost:
            t.add_row(f"[dim]{i}[/dim]",
                      f"[dim italic]{label}[/dim italic]")
        elif highlight and i in highlight:
            t.add_row(f"[bold yellow]{i}[/bold yellow]",
                      f"[bold yellow]{label}[/bold yellow]")
        else:
            t.add_row(f"[bold green]{i}[/bold green]", label)
    return t

def animated_insert(files: list, src: int, dst: int):
    t = make_table(files, highlight={src})
    console.print(t)
    h = measure_height(t)
    time.sleep(0.18)
    for _ in range(3):
        t = make_table(files)
        h = redraw_table(t, h); time.sleep(0.07)
        t = make_table(files, highlight={src})
        h = redraw_table(t, h); time.sleep(0.11)
    t = make_table(files, ghost=src)
    h = redraw_table(t, h)
    time.sleep(0.22)
    item = files.pop(src - 1)
    files.insert(dst - 1, item)
    t = make_table(files, arrived=dst)
    h = redraw_table(t, h)
    time.sleep(0.28)
    for _ in range(3):
        t = make_table(files)
        h = redraw_table(t, h); time.sleep(0.07)
        t = make_table(files, arrived=dst)
        h = redraw_table(t, h); time.sleep(0.10)
    t = make_table(files)
    redraw_table(t, h)

def _extract_pdf_ocr(path: str, page_nums: bool) -> tuple:
    try:
        import pymupdf as fitz
    except ImportError:
        import fitz
    pages = []
    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf, 1):
            items = []
            for b in page.get_text("dict")["blocks"]:
                y = b["bbox"][1]
                if b["type"] == 0:
                    txt = "\n".join("".join(s["text"] for s in ln["spans"])
                                    for ln in b["lines"])
                    if txt.strip():
                        items.append((y, 0, txt))
                else:
                    found = ocr_blob(b.get("image") or b"")
                    if found:
                        items.append((y, 1, f"[Image]\n{found}"))
            items.sort(key=lambda it: (it[0], it[1]))
            body = "\n\n".join(t for _, _, t in items)
            pages.append((f"[Page {i}]\n" if page_nums else "") + body)
    text = "\n\n".join(pages)
    return (text, "pymupdf+ocr")

def extract_pdf(path: str, page_nums: bool, ocr: bool = False) -> tuple:
    if ocr:
        try:
            return _extract_pdf_ocr(path, page_nums)
        except Exception:
            pass
    for fn, name in (
        (lambda: "\n\n".join(
            (f"[Page {i}]\n" if page_nums else "") + (p.extract_text() or '')
            for i, p in enumerate(pdfplumber.open(path).pages, 1)),
         "pdfplumber"),
        (lambda: "\n\n".join(
            (f"[Page {i}]\n" if page_nums else "") + (p.extract_text() or '')
            for i, p in enumerate(PdfReader(path).pages, 1)),
         "pypdf"),
    ):
        try:
            t = fn()
            if t.strip(): return t, name
        except Exception:
            pass
    return "[EXTRACTION ERROR]", "error"

def _para_images(p) -> list:
    blobs = []
    for blip in p._p.iter(_qn("a:blip")):
        rid = blip.get(_qn("r:embed"))
        try:
            blobs.append(p.part.rels[rid].target_part.blob)
        except Exception:
            pass
    return blobs

def extract_docx(path: str, page_nums: bool, ocr: bool = False) -> tuple:
    try:
        doc = _docx.Document(path)
        items = []
        for p in doc.paragraphs:
            t = p.text
            if ocr:
                for blob in _para_images(p):
                    found = ocr_blob(blob)
                    if found:
                        t += f"\n[Image]\n{found}"
            items.append((p, t))
        if page_nums:
            pages, cur = [], []
            for p, t in items:
                if any(run.text == "" and "w:lastRenderedPageBreak" in
                       (run._r.xml if hasattr(run, "_r") else "") for run in p.runs):
                    pages.append(cur); cur = []
                cur.append(t)
            pages.append(cur)
            text = "\n\n".join(f"[Page {i}]\n" + "\n".join(pg)
                               for i, pg in enumerate(pages, 1) if any(pg))
            return text, "python-docx"
        return "\n".join(t for _, t in items), "python-docx"
    except Exception as e:
        return f"[EXTRACTION ERROR — {e}]", "error"

def _slide_texts(shapes, ocr: bool) -> list:
    out = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            out += _slide_texts(sh.shapes, ocr)
            continue
        if ocr and sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER):
            try:
                found = ocr_blob(sh.image.blob)
                if found:
                    out.append(f"[Image]\n{found}")
            except Exception:
                pass
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                out.append(para.text)
    return out

def extract_pptx(path: str, page_nums: bool, ocr: bool = False) -> tuple:
    try:
        prs = _Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = _slide_texts(slide.shapes, ocr)
            slides.append((f"[Slide {i}]\n" if page_nums else "") + "\n".join(texts))
        return "\n\n".join(slides), "python-pptx"
    except Exception as e:
        return f"[EXTRACTION ERROR — {e}]", "error"

def extract(path: str, page_nums: bool, ocr: bool = False) -> tuple:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx": return extract_docx(path, page_nums, ocr)
    if ext == ".pptx": return extract_pptx(path, page_nums, ocr)
    return extract_pdf(path, page_nums, ocr)

def step_hdr(n: int, label: str, note: str = "") -> str:
    step = n + _STEP_OFFSET
    total = _TOTAL_STEPS
    s = (f"  [bold green]STEP {step} of {total}[/bold green]"
         f"  [dim green]╌[/dim green]  [bold white]{label}[/bold white]")
    if note:
        s += f"  [dim]{note}[/dim]"
    return s

def screen_startup():
    clear(); banner()
    w = W()
    cwd = os.getcwd()
    cwd_disp = cwd if len(cwd) < w - 18 else "…" + cwd[-(w - 19):]
    body = Text()
    body.append("\n")
    body.append("  This tool collects all  ", style="white")
    body.append("*.pdf", style="bold green")
    body.append(", ", style="white")
    body.append("*.docx", style="bold green")
    body.append("  and  ", style="white")
    body.append("*.pptx", style="bold green")
    body.append("  files in your\n", style="white")
    body.append("  current working directory and merges their text into\n", style="white")
    body.append("  a single formatted  ", style="white")
    body.append(".txt", style="bold green")
    body.append("  file.\n\n", style="white")
    body.append("  CWD ", style="bold black on green")
    body.append(f"  {cwd_disp}\n", style="bold yellow")
    body.append("\n  Nothing outside this folder is touched.\n", style="dim")
    console.print(Panel(
        body,
        title="[bold green]  contextual  [/bold green]",
        subtitle="[dim]↵ Enter to scan[/dim]",
        border_style="green",
        box=box.ROUNDED,
        width=w,
    ))
    console.print()
    console.input("  [dim]Press [bold white]Enter[/bold white] to begin …[/dim]  ")

def screen_filetype(pdfs: list, docxs: list, pptxs: list) -> list:
    clear(); banner()
    console.print(step_hdr(0, "File Type Selection"))
    console.print()
    body = Text()
    body.append("  Found  ", style="white")
    counts = []
    if pdfs: counts.append((f"{len(pdfs)} PDF",))
    if docxs: counts.append((f"{len(docxs)} DOCX",))
    if pptxs: counts.append((f"{len(pptxs)} PPTX",))
    for j, (c,) in enumerate(counts):
        if j > 0:
            body.append(",  " if j < len(counts) - 1 else "  and  ", style="white")
        body.append(c, style="bold green")
    body.append("  files in this directory.\n\n", style="white")
    options = {}
    n = 1
    if pdfs:
        body.append(f"  {n}  ", style="bold black on green")
        body.append("   PDF only   ", style="bold green")
        body.append(f"→  {len(pdfs)} file(s)\n\n", style="dim")
        options[str(n)] = pdfs; n += 1
    if docxs:
        body.append(f"  {n}  ", style="bold black on green")
        body.append("   DOCX only  ", style="bold green")
        body.append(f"→  {len(docxs)} file(s)\n\n", style="dim")
        options[str(n)] = docxs; n += 1
    if pptxs:
        body.append(f"  {n}  ", style="bold black on green")
        body.append("   PPTX only  ", style="bold green")
        body.append(f"→  {len(pptxs)} file(s)\n\n", style="dim")
        options[str(n)] = pptxs; n += 1
    all_files = sorted(pdfs + docxs + pptxs, key=numeric_key)
    body.append(f"  {n}  ", style="bold black on green")
    body.append("   All        ", style="bold green")
    body.append(f"→  {len(all_files)} file(s) interleaved by name", style="dim")
    options[str(n)] = all_files
    console.print(Panel(body, border_style="dim", box=box.ROUNDED,
                        width=W(), padding=(1, 2)))
    console.print()
    choices = " / ".join(options.keys())
    while True:
        ans = console.input(f"  [bold green]Choice >[/bold green] [dim]({choices})[/dim]  ").strip()
        if ans in options: return options[ans]
        console.print(f"  [red]Type {choices}.[/red]")

def screen_reorder(files: list) -> list:
    while True:
        clear(); banner()
        console.print(step_hdr(1, "Review & Reorder",
                               f"{len(files)} file{'s' if len(files) != 1 else ''}"))
        console.print()
        console.print(make_table(files))
        console.print()
        console.print(
            "  [dim]Type [/dim][bold green]A  B[/bold green]"
            "[dim]  to insert file [/dim][bold green]A[/bold green]"
            "[dim] at position [/dim][bold green]B[/bold green]"
            "[dim]  (others shift automatically)[/dim]\n"
            "  [dim]↵ with no input confirms and continues[/dim]"
        )
        console.print()
        raw = console.input("  [bold green]>[/bold green] ").strip()
        if not raw:
            break
        parts = raw.split()
        if len(parts) == 2 and all(p.isdigit() for p in parts):
            src, dst = int(parts[0]), int(parts[1])
            n = len(files)
            if src == dst:
                console.print("\n  [yellow]Already there — nothing to move.[/yellow]")
                time.sleep(0.8); continue
            if 1 <= src <= n and 1 <= dst <= n:
                clear(); banner()
                console.print(step_hdr(1, "Review & Reorder",
                                       f"{len(files)} files"))
                console.print()
                console.print(
                    f"  [dim]Moving [/dim][bold green]{src}[/bold green]"
                    f"[dim]  →  position [/dim][bold green]{dst}[/bold green]\n"
                )
                animated_insert(files, src, dst)
                console.print()
                console.print(
                    "  [bold green]✓[/bold green]  "
                    "[dim]Done — keep moving or ↵ to continue[/dim]"
                )
                time.sleep(1.1)
            else:
                console.print(f"\n  [red]Numbers must be 1–{n}.[/red]")
                time.sleep(0.9)
        else:
            console.print(
                "\n  [red]Enter two numbers, e.g. [bold]3 1[/bold][/red]"
            )
            time.sleep(0.9)
    return files

def screen_page_nums() -> bool:
    clear(); banner()
    w = W()
    console.print(step_hdr(2, "Output Format"))
    console.print()
    body = Text()
    body.append("  Label each extracted page in the output?\n\n", style="white")
    body.append("  y  ", style="bold black on green")
    body.append("   yes  ", style="bold green")
    body.append("→  insert [Page 1], [Page 2] … between pages\n\n", style="dim")
    body.append("  n  ", style="bold black on yellow")
    body.append("   no   ", style="bold yellow")
    body.append("→  treat each document as one continuous block of text", style="dim")
    console.print(Panel(body, border_style="dim", box=box.ROUNDED,
                        width=w, padding=(1, 2)))
    console.print()
    while True:
        ans = console.input(
            "  [bold green]Page labels?[/bold green] [dim](y / n)[/dim]  "
        ).strip().lower()
        if ans in ("y", "yes"): return True
        if ans in ("n", "no"): return False
        console.print("  [red]Type [bold]y[/bold] or [bold]n[/bold].[/red]")

def screen_ocr() -> bool:
    clear(); banner()
    w = W()
    console.print(step_hdr(3, "OCR"))
    console.print()
    body = Text()
    body.append("  OCR images inside documents?\n\n", style="white")
    body.append("  y  ", style="bold black on green")
    body.append("   yes  ", style="bold green")
    body.append("→  read text from embedded images (offline)\n\n", style="dim")
    body.append("  n  ", style="bold black on yellow")
    body.append("   no   ", style="bold yellow")
    body.append("→  skip images, text only", style="dim")
    console.print(Panel(body, border_style="dim", box=box.ROUNDED,
                        width=w, padding=(1, 2)))
    console.print()
    while True:
        ans = console.input(
            "  [bold green]OCR images inside documents?[/bold green] [dim](y / n)[/dim]  "
        ).strip().lower()
        if ans in ("y", "yes"): return True
        if ans in ("n", "no"): return False
        console.print("  [red]Type [bold]y[/bold] or [bold]n[/bold].[/red]")

def screen_filename() -> str:
    clear(); banner()
    console.print(step_hdr(4, "Output Filename"))
    console.print()
    console.print("  [dim].txt extension is added automatically[/dim]\n")
    while True:
        name = console.input("  [bold green]Filename >[/bold green] ").strip()
        if name:
            return name.removesuffix(".txt")
        console.print("  [red]Filename cannot be empty.[/red]")

def screen_extract(files: list, out_name: str, page_nums: bool, ocr: bool) -> tuple:
    clear(); banner()
    console.print(step_hdr(5, "Extracting & Writing"))
    console.print()
    start = time.time()
    results = []
    with Progress(
        SpinnerColumn(style="green"),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(bar_width=26, style="dim green", complete_style="bold green"),
        TaskProgressColumn(),
        console=console,
    ) as prog:
        task = prog.add_task("  Starting …", total=len(files))
        for f in files:
            prog.update(task, description=f"  [green]{f[:46]}[/green]")
            text, engine = extract(f, page_nums, ocr)
            prog.update(task, description=f"  [green]{f[:32]}[/green] [dim]({engine})[/dim]")
            results.append((f, text, engine))
            prog.advance(task)
    DIV = "─" * 72
    lines = [
        DIV,
        f"  CONCAT OF {len(files)} FILE{'S' if len(files) != 1 else ''}",
        f"  contextual  ·  @anshdhariwal",
        f"  Generated : {time.strftime('%Y-%m-%d  %H:%M:%S')}",
        DIV, "",
        "  TABLE OF CONTENTS", DIV,
    ]
    for i, (f, _, _) in enumerate(results, 1):
        lines.append(f"    {i:>3}.  {f}")
    lines += ["", DIV, ""]
    for i, (f, text, engine) in enumerate(results, 1):
        lines += [
            f"  [{i}/{len(results)}]  {f}",
            f"  engine : {engine}",
            DIV, "",
            text.strip() if text.strip() else "[No text extracted from this file]",
            "", DIV, "",
        ]
    out_path = f"{out_name}.txt"
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines))
    elapsed = time.time() - start
    if elapsed < 2:
        time.sleep(2 - elapsed)
    return out_path, results

def screen_done(out_path: str, results: list):
    clear(); banner()
    errors = [f for f, _, e in results if e == "error"]
    body = Text()
    body.append("\n  ✓  ", style="bold green")
    body.append(f"{len(results)}", style="bold white")
    body.append(" file(s) merged into  ", style="white")
    body.append(f"{out_path}\n", style="bold yellow")
    body.append("\n  ✓  Full path:\n     ", style="green")
    body.append(os.path.abspath(out_path), style="dim")
    if errors:
        body.append(
            f"\n\n  ⚠  {len(errors)} file(s) had extraction errors"
            " — noted inside the output\n",
            style="yellow",
        )
    else:
        body.append("\n", style="")
    console.print(Panel(body, title="[bold green]Complete[/bold green]",
                        border_style="green", box=box.ROUNDED, width=W()))
    console.print()
    time.sleep(0.5)

def screen_thanks():
    clear(); banner()
    msg = Text(justify="center")
    msg.append("glad you tried / used it bruh!\n\n", style="bold white")
    msg.append("thanks!\n\n", style="bold cyan")
    msg.append("from:  ", style="dim white")
    msg.append("maybe @anshdhariwal?", style="bold magenta")
    console.print(Panel(
        Align.center(msg),
        border_style="magenta",
        box=box.ROUNDED,
        width=W(),
        padding=(1, 8),
    ))
    time.sleep(2)
    clear()

def fast_mode(pdfs, docxs, pptxs, out_name: str, page_nums: bool, ocr: bool):
    clear(); banner()
    console.print("  [dim red]Scanning directory …[/dim red]")
    time.sleep(0.3)
    files = sorted(pdfs + docxs + pptxs, key=numeric_key)
    clear(); banner()
    console.print(
        f"  [bold red]FAST MODE[/bold red]  [dim red]╌[/dim red]"
        f"  [bold white]Detected order[/bold white]"
        f"  [dim]({len(files)} file{'s' if len(files) != 1 else ''})[/dim]\n"
    )
    console.print(make_table(files))
    console.print()
    console.print("  [dim red]No interaction in fast mode — proceeding automatically …[/dim red]")
    time.sleep(1.8)
    out_path, results = screen_extract(files, out_name, page_nums, ocr)
    screen_done(out_path, results)
    screen_thanks()

def parse_args():
    import argparse
    p = argparse.ArgumentParser(
        prog="contextual",
        description="merge every pdf, docx and pptx in the current directory into one clean txt file")
    p.add_argument("-f", "--fast", action="store_true",
                   help="no questions, merge everything: output.txt, page labels on, ocr off")
    p.add_argument("-o", "--output", metavar="NAME", help="output filename, .txt added automatically")
    p.add_argument("--pages", dest="page_nums", action="store_true", help="label pages/slides in the output")
    p.add_argument("--no-pages", dest="page_nums", action="store_false", help="keep each file as one text block")
    p.add_argument("--ocr", dest="ocr", action="store_true", help="read text out of embedded images (offline)")
    p.add_argument("--no-ocr", dest="ocr", action="store_false", help="skip embedded images entirely")
    p.set_defaults(page_nums=None, ocr=None)
    return p.parse_args()

def main():
    args = parse_args()
    global FAST_MODE, _STEP_OFFSET, _TOTAL_STEPS
    FAST_MODE = args.fast
    name = args.output.removesuffix(".txt") if args.output else None
    page_nums = args.page_nums
    use_ocr = args.ocr
    if args.fast:
        if name is None: name = "output"
        if page_nums is None: page_nums = True
        if use_ocr is None: use_ocr = False
    quiet = not sys.stdin.isatty()
    unattended = all(x is not None for x in (name, page_nums, use_ocr))

    pdfs, docxs, pptxs = discover_files()
    if not pdfs and not docxs and not pptxs:
        clear(); banner()
        console.print(Panel(
            f"[red]No [bold]*.pdf[/bold], [bold]*.docx[/bold], or [bold]*.pptx[/bold] files found in:[/red]\n\n"
            f"[dim]{os.getcwd()}[/dim]",
            title="[red]Nothing to do[/red]",
            border_style="red", box=box.ROUNDED, width=W(),
        ))
        sys.exit(1)

    if args.fast:
        fast_mode(pdfs, docxs, pptxs, name, page_nums, use_ocr)
        return

    types_present = sum(1 for t in (pdfs, docxs, pptxs) if t)
    if types_present > 1:
        if unattended:
            files = sorted(pdfs + docxs + pptxs, key=numeric_key)
        else:
            _STEP_OFFSET = 1
            _TOTAL_STEPS = 6
            files = screen_filetype(pdfs, docxs, pptxs)
    else:
        files = pdfs or docxs or pptxs

    if not quiet and not unattended:
        screen_startup()
    elif quiet and unattended:
        clear(); banner()
        console.print(f"  [dim]Running with flags — {len(files)} file(s) …[/dim]\n")
    else:
        clear(); banner()
        console.print(f"  [dim]Scanning directory …[/dim]")
        time.sleep(0.35)
        console.print(f"  [green]Found {len(files)} file(s).[/green]\n")
        time.sleep(0.3)

    files = screen_reorder(files) if not unattended or page_nums is None or use_ocr is None or name is None else files

    if page_nums is None:
        page_nums = screen_page_nums()

    if use_ocr is None:
        use_ocr = screen_ocr()

    if name is None:
        name = screen_filename()

    out_path, results = screen_extract(files, name, page_nums, use_ocr)
    screen_done(out_path, results)
    screen_thanks()

if __name__ == "__main__":
    main()
